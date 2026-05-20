import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Colors
C_DARK = RGBColor(24, 24, 27)      # zinc-900
C_ORANGE = RGBColor(249, 115, 22)  # orange-500
C_LIGHT_BG = RGBColor(244, 244, 245) # zinc-100
C_WHITE = RGBColor(255, 255, 255)
C_TEXT_DARK = RGBColor(39, 39, 42)  # zinc-800
C_TEXT_MUTED = RGBColor(113, 113, 122) # zinc-500

# Helper function to set slide background
def set_bg_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Slide 1: Cover (Dark theme)
blank_slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_slide_layout)
set_bg_color(slide1, C_DARK)

# Title & Subtitle in one text frame
txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
tf = txBox.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "AI智能食谱与拟真小票生成系统"
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = C_WHITE
p1.font.name = "Microsoft YaHei"
p1.alignment = PP_ALIGN.LEFT
p1.space_after = Pt(20)

p2 = tf.add_paragraph()
p2.text = "—— 基于大语言模型与拟物化前端架构的健康膳食系统"
p2.font.size = Pt(22)
p2.font.bold = False
p2.font.color.rgb = C_ORANGE
p2.font.name = "Microsoft YaHei"
p2.alignment = PP_ALIGN.LEFT

# Presenter Info
infoBox = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.33), Inches(1.5))
tf_info = infoBox.text_frame
tf_info.word_wrap = True

p_info = tf_info.paragraphs[0]
p_info.text = "班级学号：计科1班 202406024149\n汇报人：李相阳"
p_info.font.size = Pt(16)
p_info.font.color.rgb = RGBColor(161, 161, 170) # zinc-400
p_info.font.name = "Microsoft YaHei"
p_info.line_spacing = 1.3

# Helper to create content slide (Light theme)
def add_content_slide(title_text):
    slide = prs.slides.add_slide(blank_slide_layout)
    set_bg_color(slide, C_LIGHT_BG)
    
    # Title
    tBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.33), Inches(1.0))
    tf = tBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_DARK
    p.font.name = "Microsoft YaHei"
    
    # Divider line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.33), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_ORANGE
    shape.line.fill.background() # No border
    
    return slide

# Slide 2: 项目背景与研究现状
slide2 = add_content_slide("01 / 项目背景与用户痛点")
cBox2 = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
tf2 = cBox2.text_frame
tf2.word_wrap = True

def add_bullet(tf, title, desc, delay_pts=15):
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = "•  " + title + "："
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(4)
    
    # Add descriptions
    run = p.add_run()
    run.text = desc
    run.font.size = Pt(18)
    run.font.bold = False
    run.font.color.rgb = C_TEXT_MUTED
    p.space_after = Pt(delay_pts)

add_bullet(tf2, "日常膳食规划繁琐", "家庭烹饪缺乏针对食材库存的针对性菜谱，用户经常为“今晚吃什么”而焦虑。")
add_bullet(tf2, "健康饮食缺乏量化", "大众难以直观评估一顿饭的卡路里、营养比例、准备耗时，缺乏量化分析。")
add_bullet(tf2, "采购环节数据脱节", "从查阅菜谱到去超市采购，数据无法自动汇总。手写清单效率低且容易漏买。")
add_bullet(tf2, "本系统的核心使命", "利用大语言模型（LLM）充当 AI 营养师，打通【食材输入 -> 菜谱生成 -> 采购小票汇总 -> 移动端携带】的极简一站式闭环。")

# Slide 3: 系统架构与技术栈
slide3 = add_content_slide("02 / 系统架构与关键技术栈")
cBox3 = slide3.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
tf3 = cBox3.text_frame
tf3.word_wrap = True

add_bullet(tf3, "现代前端设计系统", "使用 React + TypeScript 确保系统类型安全，Tailwind CSS 构建极致奢华的 iOS 风格毛玻璃设计系统（ios-glass）。")
add_bullet(tf3, "多语言国际化支持 (i18n)", "集成 react-i18n，支持中英文双语一键切换，完美适配不同用户的阅读习惯。")
add_bullet(tf3, "白天/暗黑模式自适应", "采用 Tailwind CSS 暗黑模式变量，对全局输入框、文字和卡片进行了高对比度调优，杜绝背景反色带来的视感缺失。")
add_bullet(tf3, "前后端分离闭环", "前端与 Node.js 后端 API 接口流畅对接，支持用户注册登录认证、个人收藏库和云端采购清单存储。")

# Slide 4: 核心亮点 1 - AI 智能生成与冷启动交互
slide4 = add_content_slide("03 / 核心亮点：AI 智能生成与灵感交互")
cBox4 = slide4.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
tf4 = cBox4.text_frame
tf4.word_wrap = True

add_bullet(tf4, "大语言模型（LLM）精准控餐", "系统接收用户输入的配料，自动计算并排版详细的烹饪步骤、难度分级、耗时、卡路里及多维度营养分析。")
add_bullet(tf4, "“换一批”灵感搭配推荐", "在未生成内容时，提供 12 种经典家常菜配对。点击换一批，图标带有顺滑旋转微动效，一键快捷添加，解决了用户的“冷启动”输入门槛。")
add_bullet(tf4, "极致交互动效", "全站卡片均应用 iOS 弹簧缩放效果（ios-active-scale），使 Web 应用具有如同原生 iOS App 的流畅跟手感。")

# Slide 5: 核心亮点 2 - 拟真热敏小票与 Portal 传送门
slide5 = add_content_slide("04 / 核心亮点：拟真小票与全屏 Portal 传送门")
cBox5 = slide5.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
tf5 = cBox5.text_frame
tf5.word_wrap = True

add_bullet(tf5, "拟物化热敏打印小票", "模拟真实收银小票的锯齿边界、虚线排版、条形码以及打字机吐纸动效，让枯燥的数据汇总富有趣味。")
add_bullet(tf5, "React Portal 破除布局边界", "传统 Modal 容易被父容器 CSS 动画中的 transform/animation 属性干扰导致 fixed 遮罩被裁剪。我们利用 Portal 将小票完美渲染在 body 根节点下，实现 100% 全屏高清毛玻璃覆盖，消除了黑色边框 Bug。")
add_bullet(tf5, "html2canvas 精准高清导出", "重构图片截取机制，排除屏幕尺寸约束，并自动把截取克隆 DOM 下的背景、父层设为透明，输出 384px 宽的高清晰热敏小票 PNG 图，即存即走。")

# Slide 6: 项目总结与演示大纲
slide6 = add_content_slide("05 / 项目总结与演示大纲")
cBox6 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.33), Inches(5.0))
tf6 = cBox6.text_frame
tf6.word_wrap = True

add_bullet(tf6, "项目核心成效", "成功将先进的人工智能营养算法与现代拟物化美学前端交互结合，具有极高的使用价值和审美享受。")
add_bullet(tf6, "演示路线引导 (Demo Flow)", "1. 暗黑模式下AI生成 -> 2. 切换白天模式验证高对比度 -> 3. 添加采购清单 -> 4. 预览拟真小票动效 -> 5. 下载高精度小票图片。")
add_bullet(tf6, "结语", "感谢老师和各位评审！请老师批评指正。")

# Save Presentation
output_path = r"C:\Users\li\Desktop\智能食谱AI项目答辩.pptx"
prs.save(output_path)
print("SUCCESS")
